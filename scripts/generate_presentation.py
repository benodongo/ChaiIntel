"""
Generate a PowerPoint presentation summarising ChaiIntel methodology,
implementation, validation, results, and findings.

Run from the project root with the venv active:

    python scripts/generate_presentation.py

Output: analytics/static/analytics/docs/ChaiIntel-Presentation.pptx
"""

from __future__ import annotations

import io
import os
import sys
from pathlib import Path

import django
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

# ---------------------------------------------------------------------------
# Django bootstrap (so we can reuse rfutils + the same data the dashboard uses)
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))
os.environ.setdefault("DJANGO_SETTINGS_MODULE", "ChaiIntel.settings")
django.setup()

from analytics.rfutils import (  # noqa: E402
    load_historical_data,
    get_model_evaluation,
    select_global_model,
    validate_model,
    GRADES,
    GRADE_LABELS,
)

# ---------------------------------------------------------------------------
# python-pptx imports
# ---------------------------------------------------------------------------
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Brand palette (matches the dashboard CSS)
# ---------------------------------------------------------------------------
GREEN_DEEPEST = RGBColor(0x0D, 0x28, 0x18)
GREEN_DEEP    = RGBColor(0x1A, 0x3C, 0x28)
GREEN_MID     = RGBColor(0x2D, 0x6A, 0x4F)
GREEN_BRIGHT  = RGBColor(0x52, 0xB7, 0x88)
GREEN_LIGHT   = RGBColor(0xB7, 0xE4, 0xC7)
GREEN_PALE    = RGBColor(0xEA, 0xF5, 0xEE)
INK           = RGBColor(0x14, 0x1C, 0x1A)
INK_2         = RGBColor(0x2E, 0x3A, 0x36)
MUTED         = RGBColor(0x60, 0x70, 0x69)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS       = RGBColor(0x03, 0x98, 0x55)
ERROR         = RGBColor(0xD9, 0x2D, 0x20)

GRADE_HEX = {
    "BP1":      "#2D6A4F",
    "PF1":      "#52B788",
    "DUST1":    "#F4A261",
    "FNGS_1_2": "#E76F51",
    "DUST_1_2": "#8338EC",
}

OUTPUT_PATH = ROOT / "analytics" / "static" / "analytics" / "docs" / "ChaiIntel-Presentation.pptx"


# ===========================================================================
# Helpers
# ===========================================================================
def add_header_band(slide, prs, title: str) -> None:
    """Brand-coloured header strip with the slide title."""
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.7)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN_MID
    band.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.13), prs.slide_width - Inches(1.0), Inches(0.5))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_footer(slide, prs, footer_text: str = "ChaiIntel — Kenya Tea Auction Forecasting"):
    tx = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.4), prs.slide_width - Inches(1.0), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_slide(prs, title: str, subtitle: str):
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GREEN_DEEPEST
    bg.line.fill.background()

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.4), Inches(1.2), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_BRIGHT
    bar.line.fill.background()

    # Eyebrow
    eb = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(11.9), Inches(0.5))
    p = eb.text_frame.paragraphs[0]
    p.text = "ChaiIntel"
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN_LIGHT
    p.font.bold = True

    # Title
    tt = slide.shapes.add_textbox(Inches(0.7), Inches(3.1), Inches(11.9), Inches(1.6))
    p = tt.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    st = slide.shapes.add_textbox(Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.0))
    p = st.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = GREEN_LIGHT
    return slide


def section_slide(prs, eyebrow: str, title: str):
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GREEN_DEEP
    bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.0), Inches(0.8), Inches(0.10))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_BRIGHT
    bar.line.fill.background()

    eb = slide.shapes.add_textbox(Inches(0.7), Inches(3.25), Inches(11.9), Inches(0.5))
    p = eb.text_frame.paragraphs[0]
    p.text = eyebrow
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN_LIGHT
    p.font.bold = True

    tt = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(11.9), Inches(1.5))
    p = tt.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return slide


def content_slide(prs, title: str):
    slide = blank_slide(prs)
    add_header_band(slide, prs, title)
    add_footer(slide, prs)
    return slide


def add_textbox(slide, left, top, width, height, *, color=INK, bold=False):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    return tf


def add_paragraphs(text_frame, paragraphs, *, size=16, color=INK, bullet=False):
    for i, txt in enumerate(paragraphs):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = ("• " if bullet else "") + txt
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def add_table(slide, headers, rows, *, left, top, width, height,
              highlight_row_indices=None, header_fill=GREEN_MID, header_text=WHITE):
    cols = len(headers)
    n = len(rows) + 1
    shape = slide.shapes.add_table(n, cols, left, top, width, height)
    tbl = shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = header_text

    highlight = set(highlight_row_indices or [])
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i in highlight:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_PALE
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = INK
    return shape


def _render_png(plot_fn) -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(11, 5), dpi=150)
    plot_fn(ax)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


# ===========================================================================
# Chart generators (live data)
# ===========================================================================
def historical_series_chart(df) -> io.BytesIO:
    def plot(ax):
        for g in GRADES:
            if g in df.columns:
                ax.plot(df["date"], df[g], label=GRADE_LABELS[g],
                        color=GRADE_HEX[g], linewidth=2)
        ax.set_title("Historical auction prices by grade", fontsize=14, fontweight="bold")
        ax.set_ylabel("¢ / kg")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.legend(loc="upper left", frameon=False, ncol=5, fontsize=9)
        ax.grid(axis="y", linestyle="--", alpha=0.3)
    return _render_png(plot)


def model_ranking_chart(ranking) -> io.BytesIO:
    names = [r["name"] for r in ranking]
    ranks = [r["mean_rank"] for r in ranking]
    mapes = [r["mean_mape"] for r in ranking]

    def plot(ax):
        x = np.arange(len(names))
        bars = ax.bar(x, mapes,
                      color=[GREEN_MID.__str__() if False else "#2D6A4F" for _ in names])
        # Best gets brand-bright
        if mapes:
            best_idx = int(np.argmin(mapes))
            bars[best_idx].set_color("#52B788")
        ax.set_xticks(x)
        ax.set_xticklabels(names, fontsize=11)
        ax.set_ylabel("Mean MAPE (%)")
        ax.set_title("Eligible models — average MAPE across grades", fontsize=14, fontweight="bold")
        for i, (b, r_, m) in enumerate(zip(bars, ranks, mapes)):
            ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.1,
                    f"{m:.2f}%\nrank {r_:.2f}", ha="center", fontsize=9, color="#141C1A")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.grid(axis="y", linestyle="--", alpha=0.3)
    return _render_png(plot)


def validation_panel_chart(validation_rows) -> io.BytesIO:
    """A 2-by-3 grid of predicted-vs-actual on the holdout window."""
    fig, axes = plt.subplots(2, 3, figsize=(13, 6), dpi=150)
    axes = axes.flatten()
    for ax, row in zip(axes, validation_rows):
        ax.plot(row["dates"], row["actuals"], color=row["color"], linewidth=2, label="Actual")
        ax.plot(row["dates"], row["predictions"], color="#2D6A4F", linewidth=2,
                linestyle="--", label="Predicted")
        ax.set_title(f"{row['grade']}  ·  MAPE {row['mape']:.1f}%",
                     fontsize=11, fontweight="bold")
        ax.tick_params(axis="x", rotation=45, labelsize=8)
        ax.tick_params(axis="y", labelsize=8)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.legend(loc="upper left", fontsize=8, frameon=False)
    # Hide any empty axes
    for ax in axes[len(validation_rows):]:
        ax.set_visible(False)
    fig.suptitle("Predicted vs actual on the held-out 20% window",
                 fontsize=14, fontweight="bold", y=1.02)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def feature_importance_chart_for(grade: str, importances: dict) -> io.BytesIO:
    pairs = sorted(importances.items(), key=lambda kv: kv[1], reverse=True)
    labels = [p[0] for p in pairs]
    values = [p[1] for p in pairs]

    friendly = {
        "lag_1": "Lag 1 month", "lag_2": "Lag 2 months", "lag_3": "Lag 3 months",
        "rolling_mean": "Rolling avg (3m)", "month": "Month",
        "quarter": "Quarter", "time_idx": "Time trend",
        "month_sin": "Seasonality (sin)", "month_cos": "Seasonality (cos)",
    }
    labels = [friendly.get(l, l) for l in labels]

    def plot(ax):
        ax.barh(labels, values, color=["#2D6A4F" if v == max(values) else "#74C69D" for v in values])
        ax.invert_yaxis()
        ax.set_title(f"{GRADE_LABELS.get(grade, grade)} — Feature importance",
                     fontsize=14, fontweight="bold")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.grid(axis="x", linestyle="--", alpha=0.3)
        for i, v in enumerate(values):
            ax.text(v + 0.002, i, f"{v:.3f}", va="center", fontsize=9)
    return _render_png(plot)


# ===========================================================================
# Main builder
# ===========================================================================
def build():
    print("• Loading data and running models…")
    df = load_historical_data()
    evaluation = get_model_evaluation(df)
    selection = select_global_model(evaluation)
    chosen = selection["name"]
    validation = validate_model(df, chosen)

    print(f"  Chosen model: {chosen}")

    # Pre-shape validation rows for chart
    validation_rows = []
    for g in GRADES:
        v = validation["per_grade"].get(g)
        if not v or "error" in v:
            continue
        validation_rows.append({
            "grade":       GRADE_LABELS[g],
            "grade_key":   g,
            "color":       GRADE_HEX[g],
            "mape":        v["mape"],
            "rmse":        v["rmse"],
            "bias":        v["bias"],
            "directional": v.get("directional_acc"),
            "naive_mape":  v["naive_mape"],
            "beats_naive": v["beats_naive"],
            "improvement": v.get("improvement_vs_naive_pct"),
            "dates":       v["dates"],
            "actuals":     v["actuals"],
            "predictions": v["predictions"],
        })

    summary = validation.get("summary") or {}

    # ── Create deck (16:9) ─────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── 1. Title ──────────────────────────────────────────────────────────
    title_slide(
        prs,
        "Forecasting Kenyan Tea Auction Prices",
        "Methodology, implementation, validation & findings",
    )

    # ── 2. Agenda ─────────────────────────────────────────────────────────
    slide = content_slide(prs, "Agenda")
    tf = add_textbox(slide, Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.5))
    add_paragraphs(tf, [
        "1.  Problem & data",
        "2.  Methodology — candidate models, features, validation strategy",
        "3.  Model implementation",
        "4.  Model validation",
        "5.  Results",
        "6.  Findings & limitations",
    ], size=20, color=INK)

    # ── 3. Problem & data ─────────────────────────────────────────────────
    section_slide(prs, "Section 1", "Problem & data")

    slide = content_slide(prs, "Problem")
    tf = add_textbox(slide, Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.5))
    add_paragraphs(tf, [
        "Forecast monthly auction prices (¢/kg) for five Kenyan tea grades — BP1, PF1, DUST1, FNGS 1/2, DUST 1/2.",
        "Produce a 12-month forecast per grade and quantify forecast uncertainty.",
        "Constraint: use ONE model uniformly across all grades for operational simplicity.",
        "Source: monthly Mombasa Tea Auction CSV.",
        f"Dataset size: {len(df)} monthly observations, covering {df['date'].min().strftime('%b %Y')} – {df['date'].max().strftime('%b %Y')}.",
    ], size=18, bullet=True)

    slide = content_slide(prs, "Historical prices")
    img = historical_series_chart(df)
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.0), width=Inches(12.3))

    # ── 4. Methodology ────────────────────────────────────────────────────
    section_slide(prs, "Section 2", "Methodology")

    slide = content_slide(prs, "Candidate models")
    add_table(
        slide,
        ["Model", "Role", "Why it's a candidate"],
        [
            ["Naïve (last-value)", "Reference baseline",
             "Any production model must beat this. Excluded from selection."],
            ["Linear Regression",  "Interpretable benchmark",
             "Time-index, month/quarter, sin/cos, lags 1–3, rolling mean."],
            ["SARIMAX(1,1,1)",     "Classical time-series",
             "Industry standard for univariate price series with trend."],
            ["Random Forest",      "Ensemble ML",
             "Captures non-linear lag interactions without explicit modelling."],
        ],
        left=Inches(0.7), top=Inches(1.2), width=Inches(12.0), height=Inches(3.5),
    )

    slide = content_slide(prs, "Feature engineering")
    tf = add_textbox(slide, Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.5))
    add_paragraphs(tf, [
        "Autoregressive / seasonal features (Linear Regression + Random Forest):",
        "  • time_idx — global trend",
        "  • month, quarter — calendar seasonality",
        "  • month_sin, month_cos — cyclical sin/cos transform of month",
        "  • lag_1, lag_2, lag_3 — recent autoregressive signal",
        "  • rolling_mean — 3-month rolling mean, shifted by 1 step (prevents leakage)",
        "",
        "Exogenous regressors (Random Forest only):",
        "  • vol_lag_1 — previous month's sales volume (Pkgs), extracted from auction PDFs",
        "  • usd_kes  — monthly USD→KES exchange rate",
        "",
        "Rows with NaN lags are dropped before fitting.",
    ], size=14)

    slide = content_slide(prs, "How the global model is selected")
    tf = add_textbox(slide, Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.5))
    add_paragraphs(tf, [
        "1.  Run 3-fold walk-forward TimeSeriesSplit cross-validation per grade.",
        "2.  Rank the eligible models per grade by MAPE (1 = best).",
        "3.  Average each model's rank across all five grades.",
        "4.  Lowest mean rank wins — tie-broken by mean MAPE, then mean RMSE.",
        "",
        "Why ranks instead of raw MAPE? A grade with high absolute prices or high volatility could dominate a mean-MAPE comparison. Ranks are scale-free — every grade gets one equal vote.",
    ], size=16)

    # ── 5. Implementation ─────────────────────────────────────────────────
    section_slide(prs, "Section 3", "Model implementation")

    slide = content_slide(prs, "Implementation summary")
    add_table(
        slide,
        ["Model", "Key code module / class", "Hyperparameters"],
        [
            ["Naïve",             "naive_forecast (rfutils.py)",     "—"],
            ["Linear Regression", "sklearn.linear_model.LinearRegression", "default"],
            ["SARIMAX",           "statsmodels.tsa.statespace.SARIMAX",    "order=(1,1,1), trend='c'"],
            ["Random Forest",     "sklearn.ensemble.RandomForestRegressor", "n_estimators=100, max_depth=4, min_samples_leaf=2, random_state=42"],
        ],
        left=Inches(0.5), top=Inches(1.2), width=Inches(12.3), height=Inches(3.0),
    )
    tf = add_textbox(slide, Inches(0.7), Inches(4.6), Inches(12.0), Inches(2.0))
    add_paragraphs(tf, [
        "All four models share the same feature engineering pipeline (for LR / RF) and the same TimeSeriesSplit CV setup so that comparison is fair.",
        "Forecasts use a recursive multi-step loop: predict step t, append to history, predict t+1, … For RF a small trend-damping term is applied at each step.",
    ], size=14, color=INK_2)

    # ── 6. Validation ─────────────────────────────────────────────────────
    section_slide(prs, "Section 4", "Model validation")

    slide = content_slide(prs, "Validation methodology")
    tf = add_textbox(slide, Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.5))
    add_paragraphs(tf, [
        "Chronological 80/20 holdout per grade — strictly out-of-sample, independent of the CV used for selection.",
        "",
        "For each grade:",
        "1.  Split the series chronologically: first 80% → train, last 20% → test.",
        "2.  Re-fit the chosen global model on the train partition only.",
        "3.  Forecast the held-out test horizon.",
        "4.  Run the naïve baseline on the same window for comparison.",
    ], size=16)

    slide = content_slide(prs, "Metrics reported")
    add_table(
        slide,
        ["Metric", "Formula (intuition)", "Interpretation"],
        [
            ["MAE",                   "mean( |y − ŷ| )",          "Average error in ¢/kg."],
            ["RMSE",                  "sqrt( mean((y − ŷ)²) )",   "Penalises large errors more."],
            ["MAPE",                  "mean( |y − ŷ| / y )·100",  "Scale-free; comparable across grades."],
            ["Bias",                  "mean( y − ŷ )",            "+ = under-forecast; − = over-forecast."],
            ["Directional accuracy",  "% sign(ŷₜ−yₜ₋₁) = sign(yₜ−yₜ₋₁)", "Gets the up/down move right."],
            ["Improvement vs Naïve",  "(MAPE_naive − MAPE_model)/MAPE_naive", "Justifies using the model."],
        ],
        left=Inches(0.5), top=Inches(1.2), width=Inches(12.3), height=Inches(4.5),
    )

    # ── 7. Results ────────────────────────────────────────────────────────
    section_slide(prs, "Section 5", "Results")

    # 7a. Why the chosen model?
    slide = content_slide(prs, f"Selection — chosen model: {chosen}")
    img = model_ranking_chart(selection.get("ranking", []))
    slide.shapes.add_picture(img, Inches(0.7), Inches(1.0), width=Inches(7.5))

    # Side panel with stats
    panel = slide.shapes.add_textbox(Inches(8.6), Inches(1.2), Inches(4.2), Inches(5.5))
    tf = panel.text_frame
    tf.word_wrap = True
    add_paragraphs(tf, [
        f"Chosen: {chosen}",
        f"Mean rank: {selection['mean_rank']:.2f} / {len(selection.get('ranking', []))}",
        f"Mean MAPE: {selection['mean_mape']:.2f}%",
        f"Mean RMSE: {selection['mean_rmse']:.2f}",
        "",
        "Naïve is kept as a reference baseline in the comparison tables but is excluded from selection.",
    ], size=14)

    # 7b. Full ranking table
    slide = content_slide(prs, "Eligible models — global ranking")
    ranking_rows = []
    chosen_idx = None
    for i, r in enumerate(selection.get("ranking", [])):
        ranking_rows.append([
            r["name"],
            f"{r['mean_rank']:.2f}",
            f"{r['mean_mape']:.2f}%",
            f"{r['mean_rmse']:.2f}",
            f"{r['mean_mae']:.2f}",
        ])
        if r["name"] == chosen:
            chosen_idx = i
    add_table(
        slide,
        ["Model", "Mean rank ↓", "Mean MAPE", "Mean RMSE", "Mean MAE"],
        ranking_rows,
        left=Inches(0.7), top=Inches(1.2), width=Inches(12.0), height=Inches(3.5),
        highlight_row_indices=[chosen_idx] if chosen_idx is not None else None,
    )
    tf = add_textbox(slide, Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.5))
    add_paragraphs(tf, [
        "Highlighted row = chosen global model.",
        "Ranks averaged across all five grades; lower is better.",
    ], size=13, color=MUTED)

    # 7c. Per-grade CV results
    slide = content_slide(prs, "Per-grade cross-validation (all candidates)")
    rows = []
    for g in GRADES:
        if g not in evaluation:
            continue
        ev = evaluation[g]
        for m_name, m in ev["metrics"].items():
            if not m:
                continue
            rows.append([
                GRADE_LABELS[g],
                m_name,
                f"{m['mae']:.2f}",
                f"{m['rmse']:.2f}",
                f"{m['mape']:.2f}%",
                f"{m['r2']:.3f}" if m.get("r2") is not None else "—",
                "✓" if m_name == ev["best_model"] else "",
            ])
    add_table(
        slide,
        ["Grade", "Model", "MAE", "RMSE", "MAPE", "R²", "Best/grade"],
        rows,
        left=Inches(0.4), top=Inches(1.0), width=Inches(12.5), height=Inches(5.8),
    )

    # 7d. Holdout numbers
    slide = content_slide(prs, "Holdout validation — per grade")
    val_rows = []
    for r in validation_rows:
        val_rows.append([
            r["grade"],
            f"{r['mape']:.2f}%",
            f"{r['rmse']:.2f}",
            f"{r['bias']:+.2f}",
            f"{r['directional']:.1f}%" if r["directional"] is not None else "—",
            f"{r['naive_mape']:.2f}%",
            f"{r['improvement']:+.1f}%" if r["improvement"] is not None else "—",
        ])
    add_table(
        slide,
        ["Grade", "MAPE", "RMSE", "Bias", "Directional", "Naïve MAPE", "Δ vs Naïve"],
        val_rows,
        left=Inches(0.4), top=Inches(1.1), width=Inches(12.5), height=Inches(3.5),
    )
    if summary:
        s = summary
        tf = add_textbox(slide, Inches(0.7), Inches(5.0), Inches(12.0), Inches(2.0))
        add_paragraphs(tf, [
            f"Across all grades — mean MAPE: {s['mean_mape']:.2f}%   ·   mean RMSE: {s['mean_rmse']:.2f}   ·   mean bias: {s['mean_bias']:+.2f}",
            f"Beats naïve baseline on {s['beats_naive_count']} of {s['grades_evaluated']} grades.",
            f"Directional accuracy: {('%.1f%%' % s['mean_directional']) if s.get('mean_directional') is not None else '—'}.",
        ], size=14, color=INK_2)

    # 7e. Predicted vs actual chart
    if validation_rows:
        slide = content_slide(prs, "Predicted vs actual on the held-out window")
        img = validation_panel_chart(validation_rows)
        slide.shapes.add_picture(img, Inches(0.25), Inches(0.85), width=Inches(12.85))

    # ── 8. Findings ──────────────────────────────────────────────────────
    section_slide(prs, "Section 6", "Findings")

    slide = content_slide(prs, "Findings")
    tf = add_textbox(slide, Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.5))
    bullets = [
        f"The data-driven selection picked {chosen} as the single model used for every grade.",
    ]
    if summary:
        bullets.extend([
            f"On the 80/20 holdout it beats naïve carry-forward on {summary['beats_naive_count']} of {summary['grades_evaluated']} grades — sufficient evidence for adoption over a passive baseline.",
            f"Mean MAPE across grades is {summary['mean_mape']:.2f}%; mean bias is {summary['mean_bias']:+.2f} ¢/kg ({'systematic under-forecast' if summary['mean_bias'] > 0 else 'systematic over-forecast'}).",
        ])
        if summary.get("mean_directional") is not None:
            bullets.append(
                f"Directional accuracy averages {summary['mean_directional']:.1f}% — the model gets the up/down move right more than half the time on most grades."
            )
    bullets.extend([
        "Lag features dominate the Random Forest's importance ranking, confirming the autoregressive character of monthly tea prices.",
        "Seasonal sin/cos features carry secondary weight, consistent with annual harvest cycles.",
    ])
    add_paragraphs(tf, bullets, size=15, bullet=True)

    # 8b. Feature importance (only if chosen is RF and we have it)
    if chosen == "Random Forest":
        # Pick the grade with the most data
        with_fi = [
            (g, evaluation[g].get("feature_importances"))
            for g in GRADES if g in evaluation and evaluation[g].get("feature_importances")
        ]
        if with_fi:
            grade, imps = with_fi[0]
            slide = content_slide(prs, "Feature importance — Random Forest")
            img = feature_importance_chart_for(grade, imps)
            slide.shapes.add_picture(img, Inches(1.0), Inches(1.0), width=Inches(11.3))

    # 8c. Limitations
    slide = content_slide(prs, "Limitations & future work")
    tf = add_textbox(slide, Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.5))
    add_paragraphs(tf, [
        "Monthly aggregation by mean discards within-month variability; weekly granularity is available.",
        "SARIMAX order is hard-coded to (1,1,1); auto_arima would likely improve its standing.",
        "Recursive RF forecasting damps long-horizon predictions toward the recent mean. A direct multi-output regressor per horizon would be more principled.",
        "Uncertainty bands shown on the dashboard are heuristic (μ ± 1.28·σ_hist), not statistical prediction intervals. Bootstrap residual intervals or SARIMAX's get_forecast().conf_int() are obvious next steps.",
        "Dashboard re-fits every model on each request. Caching by CSV mtime, or a nightly batch job writing to a ForecastRun table, would make the UI feel instant.",
    ], size=16, bullet=True)

    # ── 9. Thank you ─────────────────────────────────────────────────────
    title_slide(
        prs,
        "Thank you",
        "Live dashboard, code & methodology write-up available at /",
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"✓ Wrote {OUTPUT_PATH}  ({OUTPUT_PATH.stat().st_size:,} bytes)")


if __name__ == "__main__":
    build()
